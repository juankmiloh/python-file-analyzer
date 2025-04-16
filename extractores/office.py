import textract
import docx
import openpyxl
from pptx import Presentation
from utils.errores import ExtractionError


def extraer_texto_doc(path):
    try:
        return textract.process(path).decode('utf-8')
    except Exception as e:
        raise ExtractionError(f"❌ Error leyendo DOC: {e}", archivo=path)

def extraer_texto_docx(path):
    try:
        doc = docx.Document(path)
        return " ".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as e:
        raise ExtractionError(f"❌ Error leyendo DOCX: {e}", archivo=path)

def extraer_texto_xlsx(path):
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
        return text
    except Exception as e:
        raise ExtractionError(f"❌ Error leyendo XLSX: {e}", archivo=path)
    
def extraer_texto_odt(path):
    try:
        return textract.process(path).decode('utf-8')
    except Exception as e:
        raise ExtractionError(f"❌ Error leyendo ODT: {e}", archivo=path)
    
def extraer_texto_pptx(path):
    try:
        prs = Presentation(path)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + " "
        return texto
    except Exception as e:
        raise ExtractionError(f"❌ Error leyendo PPTX: {e}", archivo=path)